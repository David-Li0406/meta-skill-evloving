"""
AI 技术趋势 PPT 生成脚本
风格：科技未来
总页数：15 页
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# 配图路径（brain 目录）
BRAIN_DIR = r"C:\Users\24512\.gemini\antigravity\brain\5b8c28af-1e78-4916-969f-1fceea8b51ff"

# 创建演示文稿（16:9）
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 科技未来风格配色
PRIMARY_COLOR = RGBColor(0, 229, 255)    # #00E5FF 荧光青
SECONDARY_COLOR = RGBColor(255, 23, 68)  # #FF1744 亮红
BG_COLOR = RGBColor(10, 14, 39)          # #0A0E27 深空蓝
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(176, 190, 197)

def find_image(pattern):
    """查找brain目录中的图片文件"""
    for filename in os.listdir(BRAIN_DIR):
        if pattern in filename and filename.endswith('.png'):
            return os.path.join(BRAIN_DIR, filename)
    return None

def create_cover_slide(title, subtitle):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景图
    img_path = find_image('slide_01_cover')
    if img_path:
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # 半透明遮罩
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.transparency = 0.4
    overlay.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(72)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = TEXT_WHITE
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
    tf2 = sub_box.text_frame
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(36)
    p2.runs[0].font.color.rgb = PRIMARY_COLOR

def create_content_slide(title, points, image_pattern):
    """内容页（标题 + 要点 + 图片）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    slide.shapes._spTree.remove(slide.shapes._spTree[-1])
    slide.shapes._spTree.insert(2, bg._element)
    
    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    line.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = TEXT_WHITE
    
    # 图片（左侧）
    img_path = find_image(image_pattern)
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), Inches(6.5), Inches(6))
    
    # 要点（右侧）
    content_box = slide.shapes.add_textbox(Inches(7.5), Inches(2.2), Inches(8), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {point}"
        p.space_before = Pt(14)
        p.runs[0].font.size = Pt(26)
        p.runs[0].font.color.rgb = TEXT_WHITE

# 生成所有页面
print("[进度] 正在生成第 1/15 页：封面...")
create_cover_slide("AI 的黄金时代", "2025 年人工智能技术趋势报告")

print("[进度] 正在生成第 2/15 页：目录...")
create_content_slide(
    "核心议题", 
    [
        "AI 发展简史 → 从图灵测试到通用智能",
        "2025 关键技术突破 → 多模态、推理、具身智能",
        "行业应用革命 → 医疗、金融、创意、教育",
        "技术挑战与伦理 → 安全、隐私、监管",
        "未来展望 → AGI 的下一步"
    ],
    "slide_02_contents"
)

print("[进度] 正在生成第 3/15 页：AI发展简史...")
create_content_slide(
    "从理论到现实的 75 年",
    [
        "1950 - 图灵测试：AI 哲学基础",
        "1997 - 深蓝战胜卡斯帕罗夫：算力胜利",
        "2012 - AlexNet：深度学习复兴",
        "2017 - Transformer 架构：范式转变",
        "2022 - ChatGPT：AGI 元年",
        "2025 - 多模态 AI 成为标配"
    ],
    "slide_03_history"
)

print("[进度] 正在生成第 4/15 页：Transformer架构...")
create_content_slide(
    "注意力机制改变一切",
    [
        '"Attention is All You Need" 论文开启新时代',
        "自注意力机制突破序列建模瓶颈",
        "预训练 + 微调成为主流范式",
        "GPT、BERT、T5 等模型家族爆发",
        "参数规模从亿级到万亿级跨越"
    ],
    "slide_04_transformer"
)

print("[进度] 正在生成第 5/15 页：多模态AI...")
create_content_slide(
    "视觉、语言、声音的统一",
    [
        "单一模型理解图像、文本、音频、视频",
        "Gemini、GPT-4V 引领多模态革命",
        "跨模态推理能力显著提升",
        "从 '看懂图片' 到 '理解世界'",
        "开启具身智能新篇章"
    ],
    "slide_05_multimodal"
)

print("[进度] 正在生成第 6/15 页：生成式AI...")
create_content_slide(
    "从内容消费到内容创造",
    [
        "文本生成：ChatGPT、Claude 等对话模型",
        "图像生成：DALL-E、Midjourney、Stable Diffusion",
        "视频生成：Sora、Gen-2 突破时空限制",
        "代码生成：GitHub Copilot 提升效率 10 倍",
        "音乐生成：AI 作曲进入专业领域"
    ],
    "slide_06_generative"
)

print("[进度] 正在生成第 7/15 页：推理能力...")
create_content_slide(
    "从模式识别到逻辑思考",
    [
        "Chain-of-Thought 提示技术",
        "Tree-of-Thoughts 多路径推理",
        "数学问题求解能力突破 90%",
        "科学推理辅助新发现",
        "向 AGI 的关键一步"
    ],
    "slide_07_reasoning"
)

print("[进度] 正在生成第 8/15 页：医疗应用...")
create_content_slide(
    "AI 赋能精准医疗",
    [
        "影像诊断准确率超越人类医生",
        "药物研发周期从 10 年缩短至 2 年",
        "个性化治疗方案推荐",
        "疾病预测与早期预警",
        "手术机器人辅助复杂操作"
    ],
    "slide_08_healthcare"
)

print("[进度] 正在生成第 9/15 页：金融科技...")
create_content_slide(
    "智能化重塑金融生态",
    [
        "量化交易算法日益复杂",
        "AI 风控降低坏账率 60%",
        "智能投顾普惠化",
        "反欺诈实时检测",
        "区块链 + AI 融合创新"
    ],
    "slide_09_fintech"
)

print("[进度] 正在生成第 10/15 页：创意产业...")
create_content_slide(
    "人机协作的创意新时代",
    [
        "AI 辅助文案创作效率提升 5 倍",
        "游戏 NPC 拥有真实对话能力",
        "电影特效制作成本降低 70%",
        "音乐制作民主化",
        "个性化推荐精准度突破 95%"
    ],
    "slide_10_creative"
)

print("[进度] 正在生成第 11/15 页：教育培训...")
create_content_slide(
    "个性化学习的实现",
    [
        "AI 教师 24/7 答疑解惑",
        "根据学生特点定制学习路径",
        "自动批改作业并提供改进建议",
        "沉浸式 VR/AR 教学体验",
        "知识图谱可视化学习进度"
    ],
    "slide_11_education"
)

print("[进度] 正在生成第 12/15 页：技术挑战...")
create_content_slide(
    "光明前路上的阴影",
    [
        "幻觉问题（Hallucination）仍未根除",
        "算力需求指数级增长带来能耗压力",
        "模型可解释性不足",
        "训练数据版权争议",
        "对抗样本攻击风险"
    ],
    "slide_12_challenges"
)

print("[进度] 正在生成第 13/15 页：伦理监管...")
create_content_slide(
    "在创新与安全间寻求平衡",
    [
        "AI 生成内容的真伪鉴别",
        "隐私保护与数据治理",
        "算法偏见与公平性",
        "就业替代的社会影响",
        "各国监管框架逐步成型"
    ],
    "slide_13_ethics"
)

print("[进度] 正在生成第 14/15 页：AGI未来...")
create_content_slide(
    "通用人工智能还有多远？",
    [
        "当前 AI 仍停留在 '窄智能' 阶段",
        "缺乏真正的理解和意识",
        "常识推理仍是巨大挑战",
        "乐观预测：2030-2035 年实现 AGI",
        "谨慎态度：可能需要 50 年甚至更久"
    ],
    "slide_14_agi_future"
)

print("[进度] 正在生成第 15/15 页：总结展望...")
create_content_slide(
    "拥抱 AI 时代，共创智能未来",
    [
        "AI 是工具，不是威胁",
        "人机协作将成为主流工作模式",
        "持续学习以适应技术变革",
        "关注伦理，负责任地发展 AI",
        "未来属于善用 AI 的人和组织"
    ],
    "slide_15_conclusion"
)

# 保存文件
output_path = os.path.join(BRAIN_DIR, "AI的黄金时代_2025技术趋势报告.pptx")
prs.save(output_path)

print(f"\n✅ PPT 生成完成！")
print(f"📄 文件名：AI的黄金时代_2025技术趋势报告.pptx")
print(f"📏 总页数：15 页")
print(f"🎨 风格：科技未来")
print(f"💾 大小：{os.path.getsize(output_path) / (1024*1024):.2f} MB")
print(f"📍 保存路径：{output_path}")
print(f"\n请在 WPS 或 PowerPoint 中打开进行二次编辑！")
