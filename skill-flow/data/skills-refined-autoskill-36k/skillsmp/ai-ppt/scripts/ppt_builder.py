"""
AI PPT 架构师 - PPT 生成引擎

这个脚本展示了如何使用 python-pptx 生成精美的 PPT。
实际使用时，我会根据用户的大纲动态生成完整的脚本。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import json
import os
from datetime import datetime


class PremiumPPTBuilder:
    """高端 PPT 生成器"""
    
    def __init__(self, style_path='templates/tech_future.json'):
        """初始化 PPT 生成器
        
        Args:
            style_path: 风格配置文件路径
        """
        # 加载风格配置
        with open(style_path, 'r', encoding='utf-8') as f:
            self.style = json.load(f)
        
        # 创建演示文稿（16:9 宽屏）
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        print(f"[初始化] 已加载风格：{self.style['name']}")
    
    def hex_to_rgb(self, hex_color):
        """将 HEX 颜色转换为 RGB 元组
        
        Args:
            hex_color: HEX 颜色代码（如 "#FF0000"）
            
        Returns:
            RGB 元组（如 (255, 0, 0)）
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def create_cover_slide(self, title, subtitle, image_path):
        """创建震撼的封面页
        
        布局：全屏背景图 + 半透明遮罩 + 大标题
        
        Args:
            title: 主标题
            subtitle: 副标题
            image_path: 背景图片路径
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 1. 插入全屏背景图
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                left=0, top=0,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
        
        # 2. 添加半透明深色遮罩
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0, top=0,
            width=self.prs.slide_width,
            height=self.prs.slide_height
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = self.style['design_elements']['overlay_transparency']
        overlay.line.fill.background()  # 无边框
        
        # 3. 添加标题（超大字号）
        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(3),
            width=Inches(14),
            height=Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        # 设置标题样式
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(self.style['sizes']['cover_title'])
        title_run.font.bold = True
        title_run.font.name = self.style['fonts']['title_cn']
        title_run.font.color.rgb = RGBColor(255, 255, 255)
        
        # 4. 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(5.5),
            width=Inches(14),
            height=Inches(1)
        )
        sub_frame = subtitle_box.text_frame
        sub_frame.text = subtitle
        sub_para = sub_frame.paragraphs[0]
        sub_para.alignment = PP_ALIGN.CENTER
        sub_run = sub_para.runs[0]
        sub_run.font.size = Pt(self.style['sizes']['content_title'])
        sub_run.font.name = self.style['fonts']['title_cn']
        
        # 使用主题色
        primary_color = self.hex_to_rgb(self.style['colors']['primary'])
        sub_run.font.color.rgb = RGBColor(*primary_color)
        
        print(f"[生成] 封面页：{title}")
        return slide
    
    def create_content_slide(self, title, bullet_points, image_path=None):
        """创建内容页
        
        布局：左侧图片（40%）+ 右侧文字（60%）或纯文字居中
        
        Args:
            title: 页面标题
            bullet_points: 内容要点列表
            image_path: 图片路径（可选）
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 1. 添加顶部装饰线
        line_width = self.style['design_elements']['decoration_line_width']
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0.5),
            top=Inches(0.3),
            width=Inches(3),
            height=Inches(line_width)
        )
        line.fill.solid()
        primary_rgb = self.hex_to_rgb(self.style['colors']['primary'])
        line.fill.fore_color.rgb = RGBColor(*primary_rgb)
        line.line.fill.background()
        
        # 2. 添加标题（大字号、Bold）
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.6),
            width=Inches(15),
            height=Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(self.style['sizes']['section_title'])
        title_run.font.bold = True
        title_run.font.name = self.style['fonts']['title_cn']
        
        # 根据背景色选择合适的标题颜色
        bg_color = self.style['colors']['background']
        if bg_color.startswith('#') and not bg_color.startswith('#F'):
            # 深色背景用白字
            title_run.font.color.rgb = RGBColor(255, 255, 255)
        else:
            # 浅色背景用深字
            title_run.font.color.rgb = RGBColor(33, 33, 33)
        
        # 3. 如果有图片，使用左右布局
        if image_path and os.path.exists(image_path):
            # 左侧图片
            slide.shapes.add_picture(
                image_path,
                left=Inches(0.5),
                top=Inches(2),
                width=Inches(6),
                height=Inches(6)
            )
            
            # 右侧文字区域
            content_left = Inches(7)
            content_width = Inches(8.5)
        else:
            # 无图片，文字居中
            content_left = Inches(2)
            content_width = Inches(12)
        
        # 4. 添加要点内容
        content_box = slide.shapes.add_textbox(
            left=content_left,
            top=Inches(2.5),
            width=content_width,
            height=Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {point}"
            p.level = 0
            p.space_before = Pt(12)
            
            run = p.runs[0]
            run.font.size = Pt(self.style['sizes']['body'])
            run.font.name = self.style['fonts']['body_cn']
            
            # 文字颜色
            text_color = self.hex_to_rgb(self.style['colors']['text_main'])
            run.font.color.rgb = RGBColor(*text_color)
        
        print(f"[生成] 内容页：{title}")
        return slide
    
    def save(self, output_dir='output', base_filename='演示文稿'):
        """保存 PPT 文件
        
        Args:
            output_dir: 输出目录
            base_filename: 基础文件名
            
        Returns:
            保存的文件路径
        """
        # 创建输出目录
        os.makedirs(output_dir, exist_ok=True)
        
        # 生成文件名
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{base_filename}_{timestamp}.pptx"
        filepath = os.path.join(output_dir, filename)
        
        # 保存文件
        self.prs.save(filepath)
        print(f"\n✅ PPT 生成完成：{filepath}")
        print(f"📏 总页数：{len(self.prs.slides)} 页")
        print(f"🎨 风格：{self.style['name']}")
        
        return filepath


# 示例使用
if __name__ == "__main__":
    print("=" * 50)
    print("AI PPT 架构师 - 示例脚本")
    print("=" * 50)
    
    # 创建 PPT 生成器
    builder = PremiumPPTBuilder('templates/tech_future.json')
    
    # 生成封面页
    builder.create_cover_slide(
        title="人工智能：重塑未来的力量",
        subtitle="2025 技术趋势报告",
        image_path="examples/demo_cover.png"
    )
    
    # 生成内容页 1
    builder.create_content_slide(
        title="AI 发展历程",
        bullet_points=[
            "1950: 图灵测试提出 - AI 的哲学基础",
            "2012: 深度学习复兴（AlexNet）- 视觉识别突破",
            "2017: Transformer 架构诞生 - NLP 范式转变",
            "2022: ChatGPT 引爆 AGI 浪潮 - 语言模型爆发",
            "2025: 多模态 AI 成为标配 - 视觉语言融合"
        ],
        image_path="examples/demo_content.png"
    )
    
    # 生成内容页 2
    builder.create_content_slide(
        title="关键技术突破",
        bullet_points=[
            "Transformer 架构：自注意力机制革命",
            "多模态学习：视觉与语言的统一",
            "强化学习：从AlphaGo到ChatGPT",
            "模型压缩：边缘计算时代的需求"
        ]
    )
    
    # 保存文件
    builder.save(base_filename="AI技术趋势报告")
    
    print("\n💡 提示：这只是示例脚本。")
    print("   实际使用时，我会根据你的大纲动态生成完整的 PPT。")
