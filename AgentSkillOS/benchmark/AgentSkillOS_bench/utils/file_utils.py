"""
File operation utilities.
These functions only retrieve data, they do NOT perform evaluation.
"""
import json
import csv
from pathlib import Path
from typing import Optional, List, Dict, Any
from datetime import datetime


def resolve_path(workspace: Path, path: str) -> Optional[Path]:
    """
    Resolve a file path relative to workspace, with recursive fallback.

    Resolution strategy:
    1. Try direct path: workspace / path
    2. If not found, extract the filename and recursively search the workspace

    Args:
        workspace: Base workspace directory
        path: Relative path (or just a filename) to resolve

    Returns:
        Resolved absolute Path, or None if file not found anywhere
    """
    # 1. Try direct path first
    direct = workspace / path
    if direct.exists():
        return direct

    # 2. Recursive search by filename
    filename = Path(path).name
    candidates = list(workspace.rglob(filename))
    if candidates:
        # Return the first match (shallowest path preferred)
        candidates.sort(key=lambda p: len(p.parts))
        return candidates[0]

    return None


def extract_docx_text(file_path: Path) -> Optional[str]:
    """Extract plain text from a .docx file using python-docx.

    Extracts paragraphs (newline-separated) followed by table contents
    (cells tab-separated, rows newline-separated).
    """
    try:
        from docx import Document

        doc = Document(str(file_path))
        parts: List[str] = []

        # Paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append("\t".join(cells))

        return "\n".join(parts)
    except Exception:
        return None


def extract_pdf_text(file_path: Path) -> Optional[str]:
    """Extract plain text from a PDF file using pdfplumber."""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
        return '\n'.join(texts) if texts else ''
    except Exception:
        return None


def extract_pptx_text(file_path: Path) -> Optional[str]:
    """Extract plain text from a PPTX file using python-pptx.

    Extracts text frames and table contents from all slides.
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            texts.append(para_text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append("\t".join(cells))
        return '\n'.join(texts) if texts else ''
    except Exception:
        return None


def extract_xlsx_text(file_path: Path) -> Optional[str]:
    """Extract a text representation from an .xlsx file using openpyxl.

    Loads the workbook twice — once for formulas (data_only=False) and once
    for cached computed values (data_only=True) — so that formula cells
    display as ``=SUM(A1:A10) [=150]``.
    """
    try:
        from openpyxl import load_workbook

        wb_formula = load_workbook(str(file_path), data_only=False, read_only=True)
        wb_data = load_workbook(str(file_path), data_only=True, read_only=True)
        parts: List[str] = []

        for sheet_name in wb_formula.sheetnames:
            parts.append(f"=== Sheet: {sheet_name} ===")
            ws_formula = wb_formula[sheet_name]
            ws_data = wb_data[sheet_name]

            for row_f, row_d in zip(
                ws_formula.iter_rows(max_row=200, values_only=False),
                ws_data.iter_rows(max_row=200, values_only=False),
            ):
                cells = []
                for cell_f, cell_d in zip(row_f, row_d):
                    val = cell_f.value
                    if val is None:
                        cells.append("")
                    elif isinstance(val, str) and val.startswith("="):
                        cached = cell_d.value
                        cached_str = str(cached) if cached is not None else "?"
                        cells.append(f"{val} [={cached_str}]")
                    else:
                        cells.append(str(val))
                parts.append("\t".join(cells))

        wb_formula.close()
        wb_data.close()
        return "\n".join(parts)
    except Exception:
        return None


def read_file(workspace: Path, path: str, encoding: str = "utf-8") -> Optional[str]:
    """
    Read file content as string.

    For binary document formats (.pdf, .docx, .pptx, .xlsx), extracts and
    returns plain text. For all other files, reads as text with the given
    encoding.

    Args:
        workspace: Base workspace directory
        path: Relative path to file (or filename, resolved recursively)
        encoding: File encoding (used for text files only)

    Returns:
        File content as string, or None if file doesn't exist or can't be read
    """
    file_path = resolve_path(workspace, path)
    if file_path is None:
        return None

    suffix = file_path.suffix.lower()

    if suffix == '.docx':
        return extract_docx_text(file_path)
    elif suffix == '.pdf':
        return extract_pdf_text(file_path)
    elif suffix == '.pptx':
        return extract_pptx_text(file_path)
    elif suffix == '.xlsx':
        return extract_xlsx_text(file_path)

    # Default: read as text
    try:
        with open(file_path, 'r', encoding=encoding) as f:
            return f.read()
    except Exception:
        return None


def read_json(workspace: Path, path: str) -> Optional[Dict[str, Any]]:
    """
    Read and parse JSON file.

    Args:
        workspace: Base workspace directory
        path: Relative path to JSON file

    Returns:
        Parsed JSON as dict, or None if file doesn't exist or is invalid JSON
    """
    content = read_file(workspace, path)
    if content is None:
        return None
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        return None


def read_csv(workspace: Path, path: str, encoding: str = "utf-8") -> Optional[List[Dict[str, str]]]:
    """
    Read CSV file as list of dictionaries.

    Args:
        workspace: Base workspace directory
        path: Relative path to CSV file (or filename, resolved recursively)
        encoding: File encoding

    Returns:
        List of row dictionaries, or None if file doesn't exist or can't be parsed
    """
    file_path = resolve_path(workspace, path)
    if file_path is None:
        return None
    try:
        with open(file_path, 'r', encoding=encoding) as f:
            reader = csv.DictReader(f)
            return list(reader)
    except Exception:
        return None


def check_file_exists(workspace: Path, path: str) -> bool:
    """
    Check if a file exists (searches recursively if not found at direct path).

    Args:
        workspace: Base workspace directory
        path: Relative path to file (or filename, resolved recursively)

    Returns:
        True if file exists, False otherwise
    """
    return resolve_path(workspace, path) is not None


def list_files(workspace: Path, directory: str, pattern: Optional[str] = None) -> List[str]:
    """
    List files in a directory, optionally filtered by regex pattern.

    Args:
        workspace: Base workspace directory
        directory: Relative path to directory
        pattern: Optional regex pattern to filter filenames (e.g., r"\\.png$", r"frame_.*\\.jpg")

    Returns:
        List of filenames (not full paths)
    """
    import re

    dir_path = workspace / directory
    if not dir_path.exists() or not dir_path.is_dir():
        return []

    files = [f.name for f in dir_path.iterdir() if f.is_file()]

    if pattern:
        files = [f for f in files if re.search(pattern, f)]

    return sorted(files)


def get_file_mtime(workspace: Path, path: str) -> Optional[datetime]:
    """
    Get file modification time.

    Args:
        workspace: Base workspace directory
        path: Relative path to file (or filename, resolved recursively)

    Returns:
        Modification time as datetime, or None if file doesn't exist
    """
    file_path = resolve_path(workspace, path)
    if file_path is None:
        return None
    try:
        timestamp = file_path.stat().st_mtime
        return datetime.fromtimestamp(timestamp)
    except Exception:
        return None


def get_nested_value(data: Dict[str, Any], key: str) -> Any:
    """
    Get a nested value from a dictionary using dot notation.

    Args:
        data: Dictionary to search
        key: Dot-separated key path (e.g., "result.accuracy")

    Returns:
        The value at the key path, or None if not found
    """
    current = data
    for part in key.split("."):
        if isinstance(current, dict) and part in current:
            current = current[part]
        else:
            return None
    return current
